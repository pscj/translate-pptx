from typing import List
from pptx.util import Pt

def get_english_font_size(chinese_size_pt: float) -> float:
    """Map Chinese font size to optimal English font size for visual balance.
    
    Chinese characters are denser and more complex, so English text typically
    needs to be slightly larger to achieve similar visual weight.
    """
    # Mapping table based on common PPT font sizes
    # Chinese -> English (slightly larger for readability)
    size_map = {
        8: 9,
        9: 10,
        10: 11,
        11: 12,
        12: 13,
        14: 14,   # Keep same for body text
        16: 16,   # Keep same for subtitles
        18: 18,   # Keep same for titles
        20: 20,
        22: 22,
        24: 24,
        28: 28,
        32: 32,
        36: 36,
        44: 44,
    }
    
    # Find closest match
    if chinese_size_pt in size_map:
        return size_map[chinese_size_pt]
    
    # For sizes not in map, use proportional adjustment
    # Small text (< 12pt): +1pt
    # Medium text (12-18pt): same
    # Large text (> 18pt): same
    if chinese_size_pt < 12:
        return chinese_size_pt + 1
    else:
        return chinese_size_pt

def is_chart_shape(shape):
    """Check if a shape is part of a chart or diagram."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    # Check if it's a chart type
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        return True
    # Check if shape name suggests it's part of a chart/diagram
    if hasattr(shape, 'name'):
        name_lower = shape.name.lower()
        if any(keyword in name_lower for keyword in ['chart', 'diagram', 'smartart']):
            return True
    # DO NOT treat all GROUP shapes as charts - most groups are just grouped text boxes
    # Only use RUN mode for actual charts/diagrams identified by name
    return False

def get_run_format_key(run):
    """Get a unique key representing a run's format (color, size, etc.)"""
    format_parts = []
    
    # Font size
    format_parts.append(f"size:{run.font.size}")
    
    # Font color
    if run.font.color:
        if run.font.color.type == 1:  # RGB
            format_parts.append(f"rgb:{run.font.color.rgb}")
        elif run.font.color.type == 2:  # Theme
            format_parts.append(f"theme:{run.font.color.theme_color}")
    else:
        format_parts.append("color:none")
    
    # Font style
    format_parts.append(f"bold:{run.font.bold}")
    format_parts.append(f"italic:{run.font.italic}")
    
    return "|".join(format_parts)

def group_runs_by_format(paragraph):
    """
    Group consecutive runs with the same format.
    Returns a list of groups: [(text, format_key, run_indices), ...]
    """
    if not paragraph.runs:
        return []
    
    groups = []
    current_group_text = []
    current_group_indices = []
    current_format = None
    
    for idx, run in enumerate(paragraph.runs):
        if not run.text:  # Skip empty runs
            continue
        
        run_format = get_run_format_key(run)
        
        if current_format is None or run_format == current_format:
            # Same format, add to current group
            current_group_text.append(run.text)
            current_group_indices.append(idx)
            current_format = run_format
        else:
            # Different format, save current group and start new one
            if current_group_text:
                groups.append({
                    'text': ''.join(current_group_text),
                    'format_key': current_format,
                    'run_indices': current_group_indices
                })
            # Start new group
            current_group_text = [run.text]
            current_group_indices = [idx]
            current_format = run_format
    
    # Add last group
    if current_group_text:
        groups.append({
            'text': ''.join(current_group_text),
            'format_key': current_format,
            'run_indices': current_group_indices
        })
    
    return groups

def has_multiple_formats(paragraph):
    """Check if a paragraph has multiple formats (colors, sizes, etc.)"""
    groups = group_runs_by_format(paragraph)
    return len(groups) > 1

def extract_text_from_slides(pptx_path: str):
    """Extract text from a PowerPoint presentation file and return it as a list of dicts.
    Each dict maps shape_id to extracted text data.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def extract_shape_text(shape, use_run_mode=False):
        """Extract text from a single shape, handling different shape types."""
        shape_data = []
        
        # Handle tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame"):
                        for paragraph in cell.text_frame.paragraphs:
                            # Combine all runs in a paragraph into one text
                            para_text = ''.join([run.text for run in paragraph.runs])
                            if para_text:  # Only add non-empty paragraphs
                                shape_data.append(para_text)
            return shape_data if shape_data else None
        
        # Handle grouped shapes
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                sub_data = extract_shape_text(sub_shape, use_run_mode)
                if sub_data:
                    shape_data.extend(sub_data)
            return shape_data if shape_data else None
        
        # Handle shapes with text_frame (text boxes, titles, etc.)
        elif hasattr(shape, "text_frame"):
            if use_run_mode:
                # For multi-format text: extract by format groups
                for paragraph in shape.text_frame.paragraphs:
                    groups = group_runs_by_format(paragraph)
                    for group in groups:
                        shape_data.append(group['text'])
            else:
                # For normal text: combine runs in a paragraph
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ''.join([run.text for run in paragraph.runs])
                    if para_text:  # Only add non-empty paragraphs
                        shape_data.append(para_text)
            return shape_data if shape_data else None
        
        # Handle shapes with simple text attribute
        elif hasattr(shape, "text"):
            text = shape.text.strip()
            return text if text else None
        
        return None

    prs = Presentation(pptx_path)
    all_texts = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {}  # Map shape_id to text data
        for shape_idx, shape in enumerate(slide.shapes):
            # Check if shape is a chart
            use_run_mode = is_chart_shape(shape)
            
            # Auto-detect multi-format shapes (if not already a chart)
            if not use_run_mode:
                # Check direct text_frame
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if has_multiple_formats(paragraph):
                            use_run_mode = True
                            print(f"[Extract] Slide {slide_idx}, Shape ID {shape.shape_id}: Multi-format detected, switching to RUN mode")
                            break
                # Check GROUP sub-shapes
                elif hasattr(shape, "shapes"):  # GROUP
                    for sub_shape in shape.shapes:
                        if hasattr(sub_shape, "text_frame"):
                            for paragraph in sub_shape.text_frame.paragraphs:
                                if has_multiple_formats(paragraph):
                                    use_run_mode = True
                                    print(f"[Extract] Slide {slide_idx}, Shape ID {shape.shape_id} (GROUP): Multi-format detected in sub-shape, switching to RUN mode")
                                    break
                            if use_run_mode:
                                break
                        if use_run_mode:
                            break
            
            extracted = extract_shape_text(shape, use_run_mode)
            if extracted is not None:
                shape_id = shape.shape_id
                shape_name = shape.name if hasattr(shape, 'name') else 'Unknown'
                shape_type = shape.shape_type
                mode = "RUN" if use_run_mode else "PARAGRAPH"
                print(f"[Extract] Slide {slide_idx}, Shape ID {shape_id} ({shape_name}, type={shape_type}): {mode} mode, extracted {len(extracted) if isinstance(extracted, list) else 1} items")
                slide_data[shape_id] = {
                    'text': extracted,
                    'use_run_mode': use_run_mode
                }
        all_texts.append(slide_data)

    return all_texts


def replace_text_in_slides(pptx_path: str, new_texts, output_path: str, target_language: str = "English"):
    """Replace text in a PowerPoint presentation file with new text.
    new_texts is a list of dicts, where each dict maps shape_id to translated text data.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.shapes import PP_PLACEHOLDER

    def is_title_shape(shape):
        """Check if a shape is a title placeholder."""
        try:
            if hasattr(shape, 'placeholder_format'):
                placeholder_type = shape.placeholder_format.type
                # Title, center title, or subtitle
                return placeholder_type in [PP_PLACEHOLDER.TITLE, 
                                           PP_PLACEHOLDER.CENTER_TITLE, 
                                           PP_PLACEHOLDER.SUBTITLE]
        except:
            pass
        return False

    def is_chart_shape(shape):
        """Check if a shape is part of a chart or diagram."""
        # Check if it's a chart type
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return True
        # Check if shape name suggests it's part of a chart/diagram
        if hasattr(shape, 'name'):
            name_lower = shape.name.lower()
            if any(keyword in name_lower for keyword in ['chart', 'diagram', 'smartart']):
                return True
        # DO NOT treat all GROUP shapes as charts - most groups are just grouped text boxes
        # Only use RUN mode for actual charts/diagrams identified by name
        return False

    def replace_shape_text(shape, text_data, is_title=False, use_run_mode=False):
        """Replace text in a single shape, handling different shape types."""
        para_index = 0
        run_index = 0
        
        # Handle tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    if hasattr(cell, "text_frame"):
                        for paragraph in cell.text_frame.paragraphs:
                            # Get original paragraph text
                            original_text = ''.join([run.text for run in paragraph.runs])
                            # Only process non-empty paragraphs (matching extraction logic)
                            if original_text and para_index < len(text_data):
                                translated = translate(original_text, text_data[para_index])
                                
                                if paragraph.runs:
                                    # Put translated text in first run (preserves all formatting)
                                    first_run = paragraph.runs[0]
                                    first_run.text = translated
                                    
                                    # Adjust font size for English if needed
                                    if target_language.lower() == "english" and first_run.font.size:
                                        original_size = first_run.font.size.pt
                                        first_run.font.size = Pt(get_english_font_size(original_size))
                                    
                                    # Clear text in other runs (keep runs to preserve paragraph structure)
                                    for i in range(1, len(paragraph.runs)):
                                        paragraph.runs[i].text = ""
                                # Increment index only for non-empty paragraphs
                                para_index += 1
        
        # Handle grouped shapes
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                if use_run_mode:
                    # Pass remaining data and get back the number of items consumed
                    consumed = replace_shape_text(sub_shape, text_data[run_index:], is_title, use_run_mode)
                    run_index += consumed
                else:
                    consumed = replace_shape_text(sub_shape, text_data[para_index:], is_title, use_run_mode)
                    para_index += consumed
        
        # Handle shapes with text_frame
        elif hasattr(shape, "text_frame"):
            if use_run_mode:
                # For multi-format text: replace by format groups
                for paragraph in shape.text_frame.paragraphs:
                    groups = group_runs_by_format(paragraph)
                    
                    for group in groups:
                        if run_index >= len(text_data):
                            break
                        
                        old_text = group['text']
                        new_text = text_data[run_index]
                        translated = translate(old_text, new_text)
                        
                        # Get the runs in this group
                        group_runs = [paragraph.runs[i] for i in group['run_indices']]
                        
                        # Put translated text in first run of the group
                        if group_runs:
                            first_run = group_runs[0]
                            first_run.text = translated
                            
                            # Clear text in other runs of this group
                            for run in group_runs[1:]:
                                run.text = ""
                        
                        print(f"    [GROUP mode] Replaced group {run_index}: '{old_text[:30]}...' -> '{translated[:30]}...' (format preserved, {len(group_runs)} runs)")
                        run_index += 1
            else:
                # For normal text: replace by paragraph
                for paragraph in shape.text_frame.paragraphs:
                    # Get original paragraph text
                    original_text = ''.join([run.text for run in paragraph.runs])
                    # Only process non-empty paragraphs (matching extraction logic)
                    if original_text and para_index < len(text_data):
                        translated = translate(original_text, text_data[para_index])
                        
                        if paragraph.runs:
                            # Put translated text in first run (preserves all formatting)
                            first_run = paragraph.runs[0]
                            first_run.text = translated
                            
                            # Adjust font size for English if needed
                            if target_language.lower() == "english" and first_run.font.size:
                                original_size = first_run.font.size.pt
                                first_run.font.size = Pt(get_english_font_size(original_size))
                            
                            # Clear text in other runs (keep runs to preserve paragraph structure)
                            for i in range(1, len(paragraph.runs)):
                                paragraph.runs[i].text = ""
                        # Increment index only for non-empty paragraphs
                        para_index += 1
        
        # Handle shapes with simple text attribute
        elif hasattr(shape, "text"):
            if isinstance(text_data, str):
                shape.text = translate(shape.text, text_data)
            elif para_index < len(text_data):
                shape.text = translate(shape.text, text_data[para_index])
        
        return run_index if use_run_mode else para_index

    prs = Presentation(pptx_path)

    for slide_idx, (slide, slide_data) in enumerate(zip(prs.slides, new_texts)):
        for shape in slide.shapes:
            shape_id = shape.shape_id
            
            # Check if this shape has translated text
            if shape_id in slide_data:
                shape_info = slide_data[shape_id]
                text_data = shape_info['text']
                use_run_mode = shape_info['use_run_mode']
                
                is_title = is_title_shape(shape)
                shape_name = shape.name if hasattr(shape, 'name') else 'Unknown'
                print(f"[Replace] Slide {slide_idx}, Shape ID {shape_id} ({shape_name}): {'RUN' if use_run_mode else 'PARAGRAPH'} mode, replacing {len(text_data) if isinstance(text_data, list) else 1} items")
                
                replace_shape_text(shape, text_data, is_title, use_run_mode)

    prs.save(output_path)

def translate(old_text, new_text):
    if old_text != new_text:
        print(f"Translating '{old_text}' to '{new_text}'")
    return new_text